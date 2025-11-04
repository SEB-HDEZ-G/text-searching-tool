"""
plain text & structured text extraction
file types: .txt, .csv, .json, .yaml, .xml, .md

programming & script files (treat as plain text)
.html, .py, .js, .java, .c, .cpp, .rb, .sh

office documents
file types: .docx, .xlsx, .pptx

pdf files
file type: .pdf
"""

# Imports ------------------------------------------------------------

from docx import Document
from pptx import Presentation
import openpyxl
import pdfplumber
import pytesseract
from pdf2image import convert_from_path

# Functions ------------------------------------------------------------

# Plain text and programming files text extraction
def extractPlainText(filePath):
    with open(filePath, 'r', encoding='utf-8') as file:
        extractedText = file.read()
        return extractedText
    
# Office documents text extraction
def extractOfficeText(filePath):
    if filePath.endswith('.docx'):
        doc = Document(filePath)
        text_parts = []
        
        # Extract main document paragraphs
        for para in doc.paragraphs:
            text_parts.append(para.text)
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    # Each cell contains paragraphs
                    for para in cell.paragraphs:
                        text_parts.append(para.text)
        
        # Extract headers and footers from all sections
        for section in doc.sections:
            # Header
            header = section.header
            for para in header.paragraphs:
                text_parts.append(para.text)
            
            # Footer
            footer = section.footer
            for para in footer.paragraphs:
                text_parts.append(para.text)
        
        extractedText = '\n'.join(text_parts)
        return extractedText
    
    elif filePath.endswith('.pptx'):
        prs = Presentation(filePath)
        text_parts = []
        
        for slide in prs.slides:
            # Extract text from shapes
            for shape in slide.shapes:
                # Handle text in text frames
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_parts.append(run.text)
                
                # Handle text in tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            text_parts.append(cell.text)
            
            # Extract speaker notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                text_parts.append(notes_text_frame.text)
        
        extractedText = '\n'.join(text_parts)
        return extractedText
    
    elif filePath.endswith('.xlsx'):
        wb = openpyxl.load_workbook(filePath, data_only=True)
        text_parts = []
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            
            # Extract all cell values
            for row in ws.iter_rows():
                for cell in row:
                    if cell.value is not None:
                        text_parts.append(str(cell.value))
                    
                    # Extract comments if present
                    if cell.comment:
                        text_parts.append(cell.comment.text)
        
        wb.close()
        extractedText = '\n'.join(text_parts)
        return extractedText

    else:
        raise ValueError("Unsupported office document format")
    
# PDF text extraction
def extractPDFText(filePath):
    textParts = []
    
    # Using pdfplumber for regular digital PDFs
    with pdfplumber.open(filePath) as pdf:
        for page in pdf.pages:
            # regular text
            pageText = page.extract_text()
            if pageText:
                textParts.append(pageText)
            
            # Text from tables
            tables = page.extract_tables()
            for table in tables:
                for row in table:
                    if row:
                        # Join text from non-empty cells in each row
                        rowText = ' '.join([cell for cell in row if cell])
                        textParts.append(rowText)
                        
    # Using OCR for scanned or image-based PDFs
    if not textParts:
        images = convert_from_path(filePath)
        for image in images:
            ocrText = pytesseract.image_to_string(image)
            textParts.append(ocrText)

    extractedText = ''.join(textParts)
    return extractedText

def chooseExtractionMethod(filePath):
    if filePath.endswith(('.txt', '.csv', '.json', '.yaml', '.xml', '.md',
                          '.html', '.py', '.js', '.java', '.c', '.cpp', '.rb', '.sh')):
        return extractPlainText(filePath)
    elif filePath.endswith(('.docx', '.xlsx', '.pptx')):
        return extractOfficeText(filePath)
    elif filePath.endswith('.pdf'):
        return extractPDFText(filePath)
    else:
        raise ValueError("Unsupported file format for extraction")

# Testing --------------------------------------------------------------

"""testingKey = 1  # 0-17


0: .txt, 1: .csv, 2: .json, 3: .yaml, 4: .xml, 5: .md, 6: .html, 7: .py, 
8: .js, 9: .java, 10: .c, 11: .cpp, 12: .rb, 13: .sh, 14: .docx, 
15: .xlsx, 16: .pptx, 17: .pdf

fileFormats = ['.txt', '.csv', '.json', '.yaml', '.xml', '.md',
               '.html', '.py', '.js', '.java', '.c', '.cpp', '.rb', '.sh',
               '.docx', '.xlsx', '.pptx', '.pdf']

# Extract and print the text
filePath = 'testFiles/test' + fileFormats[testingKey]
print(f"Testing extraction from: {filePath}")
print("=" * 60)

extractedText = chooseExtractionMethod(filePath)
print(extractedText)

print(extractPDFText('testFiles/test2.pdf'))"""